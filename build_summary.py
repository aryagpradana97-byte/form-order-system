#!/usr/bin/env python3
"""Build Form Order System summary PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Theme
PRIMARY = RGBColor(0x00, 0x56, 0x91)      # deep blue
ACCENT = RGBColor(0x25, 0xD3, 0x66)       # whatsapp green
DARK = RGBColor(0x0F, 0x1B, 0x2D)
GRAY = RGBColor(0x6B, 0x76, 0x86)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        for (txt, size, bold, color) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Arial"
    return tb


def title_slide():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DARK)
    add_rect(s, 0, SH - Inches(0.35), SW, Inches(0.35), ACCENT)
    add_text(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.6),
             [[("FORM ORDER SYSTEM", 44, True, WHITE)]])
    add_text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0),
             [[("Order Form  +  Order Management  +  Multi-Platform Ad Tracking", 20, False, RGBColor(0x9F, 0xB0, 0xC4))]])
    add_text(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.0),
             [[("github.com/aryagpradana97-byte/form-order-system", 14, False, ACCENT)]])


def screenshot_slide(title, img_path):
    from PIL import Image
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, DARK)
    add_text(s, Inches(0.7), Inches(0.35), Inches(12), Inches(0.8),
             [[(title, 26, True, WHITE)]])
    add_text(s, Inches(0.7), Inches(1.05), Inches(12), Inches(0.4),
             [[("Live view of the dashboard / form", 13, False, RGBColor(0x9F, 0xB0, 0xC4))]])
    im = Image.open(img_path)
    iw, ih = im.size
    max_w, max_h = Inches(12.0), Inches(5.9)
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    x = int((SW - w) / 2)
    y = int(Inches(1.55))
    s.shapes.add_picture(img_path, x, y, w, h)


def section_header(slide, text):
    add_rect(s, 0, 0, Inches(0.25), SH, PRIMARY)
    add_text(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.9),
             [[(text, 30, True, DARK)]])


def bullets(slide, items, x=Inches(0.9), y=Inches(1.7), w=Inches(11.5)):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            paras.append([("•  ", 16, True, PRIMARY), (it[0], 16, it[1], DARK)])
        else:
            paras.append([("•  ", 16, True, PRIMARY), (it, 16, False, DARK)])
    add_text(slide, x, y, w, SH - y - Inches(0.5), paras, space_after=10, line_spacing=1.05)


# ---------- Slide 1: Title ----------
title_slide()

# ---------- Slide 2: What it is ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "What It Is")
add_text(s, Inches(0.9), Inches(1.55), Inches(11.5), Inches(1.0),
         [[("A single system that runs order forms across all ad platforms and captures every "
            "lead into one place — with tracking that fires only for the submitted form.", 16, False, DARK)]],
         line_spacing=1.1)
bullets(s, [
    "Order forms for Meta, Google, TikTok, Microsoft & organic traffic",
    "One reusable tracking snippet handles ALL forms on a page",
    "Every UTM + click ID stored together (not scattered)",
    "New Google Sheet backup — never touches your existing sheets",
])

# ---------- Slide 3: Form types ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Order Form Types")
cards = [
    ("SHORT", "Name, Phone, Workshop — quick reservation"),
    ("FULL", "Name, Phone, Plate, Car, Workshop, Complaint"),
    ("DIRECT", "One-tap WhatsApp 'Consult now' button"),
    ("ORDER", "Select products/services + quantities, auto total"),
]
card_w = Inches(2.8)
gap = Inches(0.28)
x0 = Inches(0.9)
y = Inches(2.0)
for i, (t, d) in enumerate(cards):
    cx = x0 + i * (card_w + gap)
    add_rect(s, cx, y, card_w, Inches(2.6), LIGHT, line=PRIMARY)
    add_rect(s, cx, y, card_w, Inches(0.7), PRIMARY)
    add_text(s, cx + Inches(0.2), y + Inches(0.08), card_w - Inches(0.4), Inches(0.6),
             [[(t, 18, True, WHITE)]])
    add_text(s, cx + Inches(0.2), y + Inches(0.9), card_w - Inches(0.4), Inches(1.6),
             [[(d, 13, False, DARK)]], line_spacing=1.1)
add_text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.2),
         [[("Each form is embeddable separately via 'Copy Embed' — generates a complete HTML "
            "block (form + tracking + pixel IDs) ready for Elementor / WordPress / any page.",
            15, False, GRAY)]], line_spacing=1.1)

screenshot_slide("Order Form — Live Demo", "shots/demo.png")

# ---------- Slide 4: Per-form tracking ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Per-Form Tracking — the Core Fix")
bullets(s, [
    "Each form has its OWN toggle: Meta Pixel · GTM · TikTok · Google Ads",
    "Pixels fire ONLY when the form is valid AND submitted",
    "Draft / incomplete / invalid submits are still recorded as 'Draft' (not sent)",
    "Pixel ID fields per form: Meta ID, TikTok ID, GTM container, Google Ads conversion",
    "No more 'one form triggering every pixel on the page'",
])

# ---------- Slide 5: Consolidated tracking ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Consolidated Tracking (One Place)")
cols = [
    ("UTM", "utm_source, utm_medium, utm_campaign, utm_term, utm_content, utm_id"),
    ("Generic", "campaignId, adgroupId, adId, clickId, entityId"),
    ("Google", "gclid, gbraid, wbraid, gc_campaign_id, gc_adgroup_id, gc_ad_id"),
    ("Meta", "fbclid, entity_id, fb_campaign_id, fb_adset_id, fb_ad_id"),
    ("TikTok", "ttclid, tt_campaign_id, tt_adgroup_id, tt_ad_id"),
    ("Other", "msclkid · keyword · promo · site_source_name · products · total · notes"),
]
row_h = Inches(0.72)
y = Inches(1.7)
for i, (k, v) in enumerate(cols):
    add_rect(s, Inches(0.9), y + i * row_h, Inches(2.4), Inches(0.62), PRIMARY)
    add_text(s, Inches(1.05), y + i * row_h + Inches(0.08), Inches(2.2), Inches(0.5),
             [[(k, 14, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, Inches(3.3), y + i * row_h, Inches(9.1), Inches(0.62), LIGHT)
    add_text(s, Inches(3.5), y + i * row_h + Inches(0.08), Inches(8.7), Inches(0.5),
             [[(v, 12.5, False, DARK)]], anchor=MSO_ANCHOR.MIDDLE)

# ---------- Slide 6: Promos & scale ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Promos & Scale (Platform × Promo × Form)")
bullets(s, [
    ("Promos tab: create a promo per platform, then 'Generate Short + Full' forms.", True),
    ("Tracking defaults follow the platform: Meta → Meta Pixel · Google → Google Ads · TikTok → TikTok.", True),
    "Supports hundreds of combos — e.g. 10 promos × 5 platforms × 2 forms = 100",
    "All forms remain independently editable (separate tracking toggles & pixel IDs)",
])

screenshot_slide("Promos — Generate Short + Full", "shots/promos.png")

# ---------- Slide 7: Order Management ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Order Management Dashboard")
bullets(s, [
    "List all orders with search + filters (platform, form, status, date)",
    "Statuses (editable + notes/reason): Draft → Submitted → Valid → Connected → Not Connected → Deal",
    "Detail modal: full lead + tracking + products/total + events, edit status & notes",
    "'Columns' button → popup to choose which detail columns to show (Save / Cancel)",
    "Add manual order (with product selection) · Delete · Export CSV",
    "UI is in English",
])

screenshot_slide("Order Management Dashboard", "shots/orders.png")
screenshot_slide("Form Settings — Tracking Toggles & Pixel IDs", "shots/forms.png")

# ---------- Slide 8: Google Sheet mirror ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Google Sheet Mirror")
bullets(s, [
    "Ships its own Apps Script (apps-script/Code.gs)",
    "Creates a brand-new spreadsheet 'Form Order System - Leads' automatically",
    "NEVER touches your existing spreadsheet",
    "Set config.json → googleSheetWebAppUrl = your /exec URL",
    "Every submission then writes to BOTH the dashboard (JSON) and the new sheet",
])

# ---------- Slide 9: Tech / files ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Tech & Files")
bullets(s, [
    "Node.js (zero dependencies) + JSON storage + static dashboard",
    "server.js — REST API; data/ — runtime JSON (orders, forms, products, promos)",
    "public/tracker/bcs-tracker.js — embeddable tracking snippet",
    "public/index.html + app.js + styles.css — dashboard (Orders, Promos, Form Settings, Products)",
    "apps-script/Code.gs — new Apps Script → new spreadsheet",
    "Run: node server.js  →  Dashboard at http://localhost:3000",
])

screenshot_slide("Products / Services", "shots/products.png")

# ---------- Slide 10: Status & contact ----------
s = prs.slides.add_slide(BLANK)
section_header(s, "Status")
bullets(s, [
    "Verified: server boots, all API endpoints respond",
    "Form submits → recorded in dashboard (confirmed with real submissions)",
    "Google Sheet created + mirror active (existing sheet untouched)",
    "Pixel ID fields save & are used when firing",
    "Pushed to GitHub: github.com/aryagpradana97-byte/form-order-system",
])

out = "Form-Order-System-Summary.pptx"
prs.save(out)
print("Saved:", out)
